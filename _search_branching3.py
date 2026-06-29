#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
import glob, os

KEYWORD = '分支'

files = sorted(glob.glob('放射化学/*.ppt'))
files = [f for f in files if not os.path.basename(f).startswith('~')]

found_any = False
for f in files:
    try:
        prs = Presentation(f)
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text
                        if KEYWORD in text:
                            found_any = True
                            print(f"\n=== {os.path.basename(f)}, 第{i}页 ===")
                            # Print surrounding context
                            idx = text.index(KEYWORD)
                            start = max(0, idx-60)
                            end = min(len(text), idx+120)
                            print(f"  ...{text[start:end]}...")
    except Exception as e:
        print(f"Error reading {f}: {e}", file=sys.stderr)

if not found_any:
    print("在放射化学PPT中未找到「分支」")

# Also check accelerator
print(f"\n\n=== 搜索加速器PPT ===")
acc_files = sorted(glob.glob('加速器/*.ppt'))
acc_found = False
for f in acc_files:
    try:
        prs = Presentation(f)
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text
                        if KEYWORD in text:
                            acc_found = True
                            print(f"\n=== {os.path.basename(f)}, 第{i}页 ===")
                            idx = text.index(KEYWORD)
                            start = max(0, idx-60)
                            end = min(len(text), idx+120)
                            print(f"  ...{text[start:end]}...")
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)

if not acc_found:
    print("在加速器PPT中未找到「分支」")
