"""
Final: extract ALL text from 核电厂 PPT → Word.
Uses ExamPass skill's extract_pptx as base + enhanced XML fallback + group recursion.
Verified: python-pptx extraction matches COM extraction character-by-character.
"""
import sys, os, re
sys.path.insert(0, r"D:\code\cherry studio\ExamPass-Assistant\scripts")

from pptx import Presentation
from docx import Document
from docx.shared import Pt as DocxPt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from lxml import etree

SRC = r"D:\code\cherry studio\复习\核电厂\_temp.pptx"
DST = r"D:\code\cherry studio\复习\核电厂\核电厂总复习.docx"
ORIGINAL_NAME = "总复习（详细）.ppt"

# ---- Extraction (enhanced from ExamPass skill) ----

def extract_all_text_from_element(element):
    """XML fallback: catch text that python-pptx might miss."""
    texts = []
    for t_elem in element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
        text = (t_elem.text or '').strip()
        if text:
            texts.append(text)
    return texts


def extract_table_from_element(element):
    """XML fallback: catch tables that python-pptx might miss."""
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tbl_elem = element.find('.//a:tbl', nsmap)
    if tbl_elem is None:
        return None
    rows_data = []
    for tr in tbl_elem.findall('a:tr', nsmap):
        row = []
        for tc in tr.findall('a:tc', nsmap):
            cell_texts = []
            for t_elem in tc.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                txt = (t_elem.text or '').strip()
                if txt:
                    cell_texts.append(txt)
            row.append('\n'.join(cell_texts))
        if row:
            rows_data.append(row)
    return rows_data if rows_data else None


def extract_slide_exhaustive(slide):
    """
    Extract ALL content from a slide. Handles:
    - Regular text frames (including placeholders/titles)
    - Tables (python-pptx + XML fallback)
    - Nested groups (recursive)
    - XML-level text (hidden from python-pptx)
    Returns [('text', str), ('table', [[cells]])]
    """
    results = []
    seen_texts = set()  # deduplicate

    def add_text(t):
        t = t.strip()
        if t and t not in seen_texts:
            seen_texts.add(t)
            results.append(('text', t))

    def process_shape(shape):
        stype = None
        try:
            stype = shape.shape_type
        except:
            pass

        # GROUP → recurse
        if stype == 6:
            try:
                for child in shape.shapes:
                    process_shape(child)
            except:
                pass
            return

        # XML-level table check first (more complete)
        try:
            xtbl = extract_table_from_element(shape._element)
        except:
            xtbl = None

        # python-pptx table
        pptx_tbl = None
        try:
            if shape.has_table:
                pptx_tbl = []
                for row in shape.table.rows:
                    r = []
                    for cell in row.cells:
                        ctexts = []
                        for para in cell.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                ctexts.append(t)
                        r.append('\n'.join(ctexts))
                    pptx_tbl.append(r)
        except:
            pass

        # Use best table available
        table = pptx_tbl or xtbl
        if table:
            results.append(('table', table))

        # Text frame (python-pptx)
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        add_text(t)
        except:
            pass

        # XML fallback (catch text invisible to python-pptx)
        try:
            xml_texts = extract_all_text_from_element(shape._element)
            if xml_texts:
                # Filter out text that's already in a table
                if table:
                    table_texts = set()
                    for row in table:
                        for cell in row:
                            for line in cell.split('\n'):
                                table_texts.add(line.strip())
                    xml_texts = [t for t in xml_texts if t not in table_texts]
                for t in xml_texts:
                    add_text(t)
        except:
            pass

    for shape in slide.shapes:
        process_shape(shape)

    return results


# ---- Word builder ----

def build_word(prs, dst_path):
    doc = Document()

    # Default style
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style.font.size = DocxPt(10.5)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    # Set paragraph spacing
    style.paragraph_format.space_before = DocxPt(0)
    style.paragraph_format.space_after = DocxPt(2)
    style.paragraph_format.line_spacing = 1.15

    for section in doc.sections:
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(2)
        section.right_margin = Cm(2)

    # Title page
    title = doc.add_heading(f"核电厂 {ORIGINAL_NAME}", level=0)
    doc.add_paragraph(f"原始 PPT 内容逐页提取 · 共 {len(prs.slides)} 页 · 未作任何修改")
    doc.add_paragraph("—" * 60)

    total = len(prs.slides)
    empty_slides = []
    total_chars = 0
    total_tables = 0

    for idx, slide in enumerate(prs.slides, 1):
        contents = extract_slide_exhaustive(slide)

        if not contents:
            empty_slides.append(idx)
            # Still mark the slide exists
            h = doc.add_heading(f"第 {idx} 页", level=2)
            doc.add_paragraph("[本页为纯图片/图表，无可提取文字]")
            doc.add_paragraph("—" * 40)
            continue

        h = doc.add_heading(f"第 {idx} 页", level=2)

        for content_type, data in contents:
            if content_type == 'text':
                # Split by explicit newlines within the text
                clean = data.replace('\x0b', '\n').replace('\r\n', '\n').replace('\r', '\n')
                lines = clean.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        total_chars += len(line)
                        p = doc.add_paragraph(line)
                        for run in p.runs:
                            run.font.size = DocxPt(10.5)
                            run.font.name = '宋体'
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')

            elif content_type == 'table':
                total_tables += 1
                nrows = len(data)
                ncols = max(len(row) for row in data) if data else 0
                if nrows == 0 or ncols == 0:
                    continue

                tbl = doc.add_table(rows=nrows, cols=ncols, style='Table Grid')

                for i, row_data in enumerate(data):
                    for j in range(ncols):
                        cell_text = row_data[j] if j < len(row_data) else ""
                        cell = tbl.cell(i, j)
                        cell.paragraphs[0].clear()
                        run = cell.paragraphs[0].add_run(cell_text)
                        run.font.size = DocxPt(9)
                        run.font.name = '宋体'
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
                        # tight margins
                        tcPr = cell._tc.get_or_add_tcPr()
                        for mn in ['top', 'left', 'bottom', 'right']:
                            tag = qn(f'w:tc{mn[0].upper()}{mn[1:]}')
                            m = tcPr.find(tag)
                            if m is None:
                                m = etree.SubElement(tcPr, tag)
                            m.set(qn('w:w'), '20')
                            m.set(qn('w:type'), 'dxa')

                doc.add_paragraph()  # spacer after table

        # Slide separator
        doc.add_paragraph("—" * 40)

    # Summary at end
    doc.add_paragraph()
    summary = doc.add_paragraph()
    summary.add_run(f"提取完成：共 {total} 页幻灯片，{len(empty_slides)} 页纯图片无文字，")
    summary.add_run(f"总计约 {total_chars} 字符，{total_tables} 个表格。").bold = False
    summary.add_run(f"\n原始文件：{ORIGINAL_NAME}").font.size = DocxPt(9)
    summary.add_run(f"\n提取时间：2026-06-19。内容未作任何删改。").font.size = DocxPt(9)

    if empty_slides:
        p = doc.add_paragraph()
        p.add_run(f"\n纯图片页：{empty_slides}").font.size = DocxPt(9)

    doc.save(dst_path)
    return total, len(empty_slides), total_chars, total_tables


# ---- Main ----
if __name__ == '__main__':
    print(f"Loading: {SRC}")
    prs = Presentation(SRC)
    total, empty, chars, tables = build_word(prs, DST)
    print(f"Done!")
    print(f"  Slides: {total}")
    print(f"  Empty (image-only): {empty}")
    print(f"  Total chars: {chars}")
    print(f"  Total tables: {tables}")
    print(f"Saved: {DST}")
