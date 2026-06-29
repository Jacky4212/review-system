#!/usr/bin/env python3
"""Extract text from all PPTX files in the current directory."""
import os
from pptx import Presentation

out_dir = "提取内容"
os.makedirs(out_dir, exist_ok=True)

files = sorted([f for f in os.listdir('.') if f.endswith('.pptx') and not f.startswith('~')])

# Skip (2) duplicates if originals exist
filtered = []
for fname in files:
    if '(2)' in fname and fname.replace('(2)', '') in files:
        continue
    filtered.append(fname)

for fname in filtered:
    print("\n" + "="*60)
    print("=== 正在提取: %s" % fname)

    try:
        prs = Presentation(fname)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(cells))

            if slide_text:
                slides_text.append("--- 第%d页 ---\n%s" % (i, "\n".join(slide_text)))

        out_name = fname.replace('.pptx', '.txt')
        out_path = os.path.join(out_dir, out_name)
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write("来源: %s\n" % fname)
            f.write("总页数: %d\n\n" % len(prs.slides))
            f.write("\n\n".join(slides_text))

        print("=== 已保存 (%d 页有文字)" % len(slides_text))

    except Exception as e:
        print("=== 提取失败: %s" % e)

print("\n\n=== 全部提取完成！")
