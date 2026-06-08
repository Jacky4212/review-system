"""
检查原始PPT与提取文本的一致性
"""
import os
from pptx import Presentation, util as pptx_util

ppt_dir = r"D:\code\cherry studio\复习\放射化学"
out_dir = r"D:\code\cherry studio\复习\output"

# Also check slide-level detail for ch1 and ch2
for fname in sorted(os.listdir(ppt_dir)):
    if not fname.endswith('.ppt') or fname.startswith('~$'):
        continue
    path = os.path.join(ppt_dir, fname)
    try:
        prs = Presentation(path)
        slides = list(prs.slides)
    except Exception as e:
        print(f"Cannot read: {fname}: {e}")
        continue

    # Find corresponding txt file
    base = os.path.splitext(fname)[0]
    txt_path = os.path.join(out_dir, base + '.txt')
    if not os.path.exists(txt_path):
        txt_path = os.path.join(out_dir, fname.replace('.ppt', '.txt'))
    if not os.path.exists(txt_path):
        print(f"NO TXT for {fname}: {base}")
        continue

    with open(txt_path, 'r', encoding='utf-8') as f:
        txt_content = f.read()

    txt_slides = txt_content.count("第 ") - txt_content.count("第 1 页")

    # Extract ppt content slide by slide
    ppt_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text.strip())
        ppt_text.append("\n".join(slide_texts))

    print(f"\n=== {fname} ===")
    print(f"  PPT slides: {len(slides)}, TXT slides: {txt_slides}")

    # Report slide-by-slide for ch1 and ch2
    if "第一章" in fname or "第二章" in fname:
        for i, (slide, stext) in enumerate(zip(slides, ppt_text)):
            # Find this slide in txt
            marker = f"===== 第 {i+1} 页 ====="
            if marker not in txt_content:
                print(f"  SLIDE {i+1}: MISSING from text!")
                continue
            # Check text length
            txt_part = txt_content.split(marker)
            if len(txt_part) > 1:
                txt_slide = txt_part[1].split("===== 第")[0] if "===== 第" in txt_part[1] else txt_part[1]
                ppt_len = len(stext)
                txt_len = len(txt_slide.strip())
                if abs(ppt_len - txt_len) > 50:
                    print(f"  Slide {i+1}: ppt={ppt_len}c, txt={txt_len}c {('(big diff!)' if abs(ppt_len-txt_len)>200 else '')}")
